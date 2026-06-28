#!/usr/bin/env python
"""pptx から文字を取り除いたスライドを画像 (PNG) として書き出すスクリプト。

各スライドのテキスト（テキストフレーム・表セル・グループ内図形を含む）を空に
してから、残ったオブジェクト（図・図形・画像）だけのスライドを 1 枚ずつ PNG に
レンダリングする。

パイプライン:
  python-pptx でテキスト除去 → 一時 pptx 保存 → soffice で PDF 化
  → PyMuPDF (fitz) で各ページを PNG レンダリング

実行には conda の py313 環境と soffice が必要:
    /opt/conda/envs/py313/bin/python src/pptx_to_md/render_images.py --all

使い方:
    python src/pptx_to_md/render_images.py --all
    python src/pptx_to_md/render_images.py pptx/foo.pptx --dpi 200

結果は 1 ファイル 1 行の JSON で stdout に出力する。
"""
from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# LibreOffice 生成 PDF の構造ツリーに対する非致命的な MuPDF 警告を抑制する。
fitz.TOOLS.mupdf_display_errors(False)

SOFFICE = "soffice"

# スライドマスター/レイアウト上で「装飾オブジェクト」とみなして削除する図形要素。
_SHAPE_TAGS = {
    qn("p:sp"),
    qn("p:pic"),
    qn("p:graphicFrame"),
    qn("p:grpSp"),
    qn("p:cxnSp"),
    qn("p:contentPart"),
}


def _strip_text_frame(text_frame) -> None:
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.text = ""


def strip_text_from_shape(shape) -> None:
    """1 つの図形（グループ・表を含む）から再帰的にテキストを取り除く。"""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            strip_text_from_shape(child)
        return
    if shape.has_text_frame:
        _strip_text_frame(shape.text_frame)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _strip_text_frame(cell.text_frame)


def remove_tables_from_shapes(shapes) -> None:
    """図形コレクションから表 (table) を再帰的に削除する。

    表はスライド本体に直接置かれた graphicFrame オブジェクトなので、spTree から
    要素ごと取り除く。グループ内の表も対象にする。
    """
    for shape in list(shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            remove_tables_from_shapes(shape.shapes)
        elif getattr(shape, "has_table", False):
            shape._element.getparent().remove(shape._element)


def _clear_template_part(part) -> None:
    """スライドマスター / レイアウトから装飾図形と背景を取り除く。

    プレースホルダを含む全図形を spTree から削除し、背景 (p:bg) も外す。
    これによりレンダリング時にテンプレート由来のオブジェクト（バナー・ロゴ・
    ページ番号の枠など）が描画されなくなる。スライド本体に直接置かれた図形は
    別の spTree にあるため影響を受けない。
    """
    sp_tree = part.shapes._spTree
    for child in list(sp_tree):
        if child.tag in _SHAPE_TAGS:
            sp_tree.remove(child)
    csld = part._element.find(qn("p:cSld"))
    if csld is not None:
        bg = csld.find(qn("p:bg"))
        if bg is not None:
            csld.remove(bg)


def drop_template_objects(prs) -> None:
    """全スライドマスターとレイアウトからテンプレート由来の装飾を削除する。"""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            _clear_template_part(layout)
        _clear_template_part(master)


def shape_is_content(shape) -> bool:
    """その図形が「描画されるコンテンツ」かどうかを判定する。

    - 画像 (picture): コンテンツ（プレースホルダ画像も含む）
    - グループ: 子に 1 つでもコンテンツがあればコンテンツ
    - 表 (table): 残っていれば（keep-tables 時）コンテンツ
    - プレースホルダ（タイトル枠などのテンプレート構造）: 中身が空なら非コンテンツ
    - それ以外の図形（作者が直接置いた図形・矢印・テキストボックス等）: コンテンツ
    """
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return any(shape_is_content(child) for child in shape.shapes)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if getattr(shape, "has_table", False):
        return True
    if shape.is_placeholder:
        return False
    return True


def slide_has_content(slide) -> bool:
    """スライド本体に描画されるコンテンツ図形が 1 つでもあれば True。"""
    return any(shape_is_content(shape) for shape in slide.shapes)


def make_textless_pptx(
    src: Path, dst: Path, drop_template: bool = True, drop_tables: bool = True
) -> list[int]:
    """src からテキストを除いた pptx を dst に保存する。

    drop_template=True のとき、スライドマスター/レイアウト由来のテンプレート
    オブジェクト（背景・バナー・ロゴ等）も併せて取り除く。
    drop_tables=True のとき、スライド本体の表 (table) も取り除く。

    戻り値は、除去後にコンテンツ図形が何も残らなかったスライド番号（1 始まり）の
    一覧。呼び出し側はこれをレンダリングからスキップできる。
    """
    prs = Presentation(str(src))
    empty_slides: list[int] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            strip_text_from_shape(shape)
        if drop_tables:
            remove_tables_from_shapes(slide.shapes)
        if not slide_has_content(slide):
            empty_slides.append(idx)
    if drop_template:
        drop_template_objects(prs)
    prs.save(str(dst))
    return empty_slides


def pptx_to_pdf(pptx_path: Path, outdir: Path, profile_dir: Path) -> Path:
    """soffice で pptx を PDF に変換し、生成された PDF のパスを返す。"""
    subprocess.run(
        [
            SOFFICE,
            "--headless",
            f"-env:UserInstallation=file://{profile_dir}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(outdir),
            str(pptx_path),
        ],
        check=True,
        capture_output=True,
    )
    pdf_path = outdir / (pptx_path.stem + ".pdf")
    if not pdf_path.is_file():
        raise RuntimeError("soffice did not produce a PDF")
    return pdf_path


def pdf_to_pngs(
    pdf_path: Path, outdir: Path, stem: str, dpi: int, skip_slides: set[int] | None = None
) -> list[str]:
    """PDF の各ページを PNG にレンダリングし、保存した画像パスの一覧を返す。

    skip_slides に含まれるページ（スライド番号, 1 始まり）は保存しない。
    ファイル名にはスライド番号を埋め込むため、スキップ分は番号が飛ぶ。
    """
    skip_slides = skip_slides or set()
    outdir.mkdir(parents=True, exist_ok=True)
    # 再生成時に、前回の（スキップ対象になった分も含む）古い画像が残らないよう掃除する。
    for stale in outdir.glob(f"{stem}_slide_*.png"):
        stale.unlink()
    images: list[str] = []
    with fitz.open(str(pdf_path)) as doc:
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        width = len(str(doc.page_count))
        for i, page in enumerate(doc, start=1):
            if i in skip_slides:
                continue
            pix = page.get_pixmap(matrix=matrix)
            out = outdir / f"{stem}_slide_{i:0{width}d}.png"
            pix.save(str(out))
            images.append(str(out))
    return images


def render_one(
    pptx_path: Path,
    outroot: Path,
    dpi: int,
    drop_template: bool = True,
    drop_tables: bool = True,
) -> dict:
    result: dict = {
        "input": str(pptx_path),
        "outdir": None,
        "images": [],
        "count": 0,
        "skipped": [],
        "ok": False,
    }
    if not pptx_path.is_file():
        result["error"] = "input file not found"
        return result

    img_dir = outroot / pptx_path.stem / "slide_images"
    # 一時ファイルは /tmp を使わず、出力ルート配下に作る（CLAUDE.md 規約）。
    work = Path(tempfile.mkdtemp(prefix=".render_", dir=str(outroot)))
    try:
        textless = work / (pptx_path.stem + "_notext.pptx")
        empty_slides = make_textless_pptx(
            pptx_path, textless, drop_template=drop_template, drop_tables=drop_tables
        )
        profile = work / "lo_profile"
        pdf = pptx_to_pdf(textless, work, profile)
        images = pdf_to_pngs(
            pdf, img_dir, pptx_path.stem, dpi, skip_slides=set(empty_slides)
        )
        result.update(
            outdir=str(img_dir),
            images=images,
            count=len(images),
            skipped=empty_slides,
            ok=True,
        )
    except subprocess.CalledProcessError as exc:
        result["error"] = f"soffice failed: {exc.stderr.decode('utf-8', 'replace')[:300]}"
    except Exception as exc:
        result["error"] = f"render failed: {exc}"
    finally:
        shutil.rmtree(work, ignore_errors=True)
    return result


def collect_inputs(args: argparse.Namespace) -> list[Path]:
    if args.all:
        return sorted(Path(args.pptx_dir).glob("**/*.pptx"))
    return [Path(p) for p in args.inputs]


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="pptx から文字を除いたスライドを PNG 化する")
    parser.add_argument("inputs", nargs="*", help="対象の .pptx ファイル")
    parser.add_argument("--all", action="store_true", help="--pptx-dir 配下の全 .pptx を対象にする")
    parser.add_argument("--pptx-dir", default="pptx", help="--all のときの探索元 (既定: pptx)")
    parser.add_argument("--outdir", default="docs_from_agent", help="出力ルート (既定: docs_from_agent)")
    parser.add_argument("--dpi", type=int, default=150, help="レンダリング解像度 (既定: 150)")
    parser.add_argument(
        "--keep-template",
        dest="drop_template",
        action="store_false",
        default=True,
        help="テンプレート（マスター/レイアウト）由来の装飾・背景を残す。既定では削除する。",
    )
    parser.add_argument(
        "--keep-tables",
        dest="drop_tables",
        action="store_false",
        default=True,
        help="スライド本体の表 (table) を残す。既定では削除する。",
    )
    args = parser.parse_args(argv)

    inputs = collect_inputs(args)
    if not inputs:
        print("no .pptx files found", file=sys.stderr)
        return 1

    outroot = Path(args.outdir)
    outroot.mkdir(parents=True, exist_ok=True)
    failed = 0
    for pptx_path in inputs:
        res = render_one(
            pptx_path,
            outroot,
            args.dpi,
            drop_template=args.drop_template,
            drop_tables=args.drop_tables,
        )
        if not res["ok"]:
            failed += 1
        print(json.dumps(res, ensure_ascii=False))

    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
