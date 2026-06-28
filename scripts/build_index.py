#!/usr/bin/env python
"""パワーポイントごとに、各スライドの成果物をまとめた index.json を生成する。

`docs_from_agent/<名前>/` 配下の slide_images / images / tables を走査し、
ファイル名に埋め込まれたスライド番号 (`slide_NNN` / `slideNNN`) を手がかりに
スライド単位へ集約する。スライド数・タイトルは元 pptx から取得する。

出力: `docs_from_agent/<名前>/index.json`（機械処理しやすい JSON）

実行には py313 環境 (python-pptx) が必要:
    /opt/conda/envs/py313/bin/python src/pptx_to_md/build_index.py --all

使い方:
    python src/pptx_to_md/build_index.py --all
    python src/pptx_to_md/build_index.py pptx/foo.pptx

結果は 1 ファイル 1 行の JSON で stdout に出力する。
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation

# ファイル名からスライド番号を取り出す（slide_004 / slide004 のどちらも対応）。
_SLIDE_RE = re.compile(r"_slide_?(\d+)")
# images / tables の連番 M を取り出す（並び順用）。
_SEQ_RE = re.compile(r"_(?:img|table)(\d+)")
# md をスライド境界で分割する markitdown のコメント。
_SLIDE_COMMENT_RE = re.compile(r"<!--\s*Slide number:\s*(\d+)\s*-->")
# md 本文中の画像参照（テキストからは除き、images フィールドで持つ）。
_IMG_REF_RE = re.compile(r"!\[[^\]]*\]\([^)]*\)\s*")
# 発表者ノートの区切り。
_NOTES_RE = re.compile(r"\n###\s*Notes:\s*\n?")
# 先頭見出し行（# タイトル）。
_HEADING_RE = re.compile(r"^#\s+(.*)$", re.MULTILINE)


def parse_md_slides(md_path: Path) -> dict[int, dict]:
    """生成済み md をスライド境界で分割し、{番号: {text, notes, md_title}} を返す。

    text は本文（画像参照を除去、表の markdown は保持）、notes は発表者ノート。
    """
    if not md_path.is_file():
        return {}
    raw = md_path.read_text(encoding="utf-8")
    parts = _SLIDE_COMMENT_RE.split(raw)
    # split 結果: [先頭, num1, block1, num2, block2, ...]
    out: dict[int, dict] = {}
    for i in range(1, len(parts), 2):
        num = int(parts[i])
        block = parts[i + 1] if i + 1 < len(parts) else ""
        body, notes = block, ""
        m = _NOTES_RE.search(block)
        if m:
            body, notes = block[: m.start()], block[m.end():]
        heading = _HEADING_RE.search(body)
        md_title = heading.group(1).strip() if heading else ""
        text = _IMG_REF_RE.sub("", body).strip()
        out[num] = {"text": text, "notes": notes.strip(), "md_title": md_title}
    return out


def _first_line(text: str) -> str:
    for line in text.splitlines():
        line = line.lstrip("# ").strip()
        if line:
            return line
    return ""


def _slide_no(filename: str) -> int | None:
    m = _SLIDE_RE.search(filename)
    return int(m.group(1)) if m else None


def _seq_no(filename: str) -> int:
    m = _SEQ_RE.search(filename)
    return int(m.group(1)) if m else 0


def _scan(subdir: Path, rel_prefix: str) -> dict[int, list[str]]:
    """subdir 内のファイルをスライド番号ごとに集めて返す（相対パス、連番順）。"""
    by_slide: dict[int, list[str]] = {}
    if not subdir.is_dir():
        return by_slide
    files = [p for p in subdir.iterdir() if p.is_file()]
    for p in sorted(files, key=lambda p: (_slide_no(p.name) or 0, _seq_no(p.name), p.name)):
        n = _slide_no(p.name)
        if n is None:
            continue
        by_slide.setdefault(n, []).append(f"{rel_prefix}/{p.name}")
    return by_slide


def _slide_titles(prs) -> dict[int, str]:
    titles: dict[int, str] = {}
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text:
                # 改行・垂直タブ等の空白類を 1 つの空白に正規化する。
                title = re.sub(r"\s+", " ", slide.shapes.title.text).strip()
        except Exception:
            pass
        titles[idx] = title
    return titles


def build_one(pptx_path: Path, outroot: Path) -> dict:
    result: dict = {"input": str(pptx_path), "index": None, "slides": None, "ok": False}
    if not pptx_path.is_file():
        result["error"] = "input file not found"
        return result

    stem = pptx_path.stem
    base = outroot / stem
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        result["error"] = f"open failed: {exc}"
        return result

    slide_count = len(prs.slides)
    titles = _slide_titles(prs)

    slide_images = _scan(base / "slide_images", "slide_images")
    images = _scan(base / "images", "images")
    tables = _scan(base / "tables", "tables")

    md_name = f"{stem}.md"
    markdown = md_name if (base / md_name).is_file() else None
    md_slides = parse_md_slides(base / md_name)

    slides = []
    for n in range(1, slide_count + 1):
        si = slide_images.get(n, [])
        md = md_slides.get(n, {})
        text = md.get("text", "")
        # タイトルは pptx の title 枠 → md の見出し → 本文 1 行目 の順に補完する。
        title = titles.get(n, "") or md.get("md_title", "") or _first_line(text)
        slides.append(
            {
                "id": f"{stem}#{n:03d}",
                "slide": n,
                "title": title,
                "text": text,
                "notes": md.get("notes", ""),
                # text が空でもメディアがあれば情報はあるが、has_content はテキスト有無を指す。
                "has_content": bool(text),
                # slide_image は基本 1 枚（空スライドはスキップされ無い場合あり）。
                "slide_image": si[0] if si else None,
                "images": images.get(n, []),
                "tables": tables.get(n, []),
            }
        )

    deck_title = titles.get(1, "") or (slides[0]["title"] if slides else "") or stem
    index = {
        "name": stem,
        "title": deck_title,
        "source_pptx": str(pptx_path),
        "slide_count": slide_count,
        "markdown": markdown,
        "totals": {
            "slide_images": sum(len(v) for v in slide_images.values()),
            "images": sum(len(v) for v in images.values()),
            "tables": sum(len(v) for v in tables.values()),
        },
        "slides": slides,
    }

    base.mkdir(parents=True, exist_ok=True)
    out_path = base / "index.json"
    out_path.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")

    result.update(index=str(out_path), slides=slide_count, ok=True)
    return result


def collect_inputs(args: argparse.Namespace) -> list[Path]:
    if args.all:
        return sorted(Path(args.pptx_dir).glob("**/*.pptx"))
    return [Path(p) for p in args.inputs]


def _prefix_path(name: str, p: str | None) -> str | None:
    """deck 相対パスを outroot 相対（<名前>/...）に直す。"""
    return f"{name}/{p}" if p else None


def build_chunks_jsonl(outroot: Path) -> tuple[Path, int]:
    """全 <名前>/index.json を走査し、1 行 1 スライドの chunks.jsonl を生成する。

    各行はそのスライド単体で RAG のチャンクになる（本文・メタデータ・メディアパス）。
    メディアパスは outroot からの相対（<名前>/...）。
    """
    out_path = outroot / "chunks.jsonl"
    count = 0
    with out_path.open("w", encoding="utf-8") as f:
        for index_path in sorted(outroot.glob("*/index.json")):
            data = json.loads(index_path.read_text(encoding="utf-8"))
            name = data["name"]
            for s in data["slides"]:
                chunk = {
                    "id": s["id"],
                    "deck": name,
                    "deck_title": data.get("title", ""),
                    "source_pptx": data.get("source_pptx", ""),
                    "slide": s["slide"],
                    "slide_count": data.get("slide_count"),
                    "title": s["title"],
                    "text": s["text"],
                    "notes": s["notes"],
                    "has_content": s["has_content"],
                    "slide_image": _prefix_path(name, s["slide_image"]),
                    "images": [_prefix_path(name, p) for p in s["images"]],
                    "tables": [_prefix_path(name, p) for p in s["tables"]],
                }
                f.write(json.dumps(chunk, ensure_ascii=False) + "\n")
                count += 1
    return out_path, count


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="pptx ごとの index.json を生成する")
    parser.add_argument("inputs", nargs="*", help="対象の .pptx ファイル")
    parser.add_argument("--all", action="store_true", help="--pptx-dir 配下の全 .pptx を対象にする")
    parser.add_argument("--pptx-dir", default="pptx", help="--all のときの探索元 (既定: pptx)")
    parser.add_argument("--outdir", default="docs_from_agent", help="出力ルート (既定: docs_from_agent)")
    parser.add_argument("--no-chunks", action="store_true", help="横断コーパス chunks.jsonl を生成しない")
    args = parser.parse_args(argv)

    inputs = collect_inputs(args)
    if not inputs:
        print("no .pptx files found", file=sys.stderr)
        return 1

    outroot = Path(args.outdir)
    failed = 0
    for pptx_path in inputs:
        res = build_one(pptx_path, outroot)
        if not res["ok"]:
            failed += 1
        print(json.dumps(res, ensure_ascii=False))

    # 全 index.json から横断コーパス chunks.jsonl を（常に全体で）再生成する。
    if not args.no_chunks:
        chunks_path, n = build_chunks_jsonl(outroot)
        print(json.dumps({"chunks": str(chunks_path), "count": n, "ok": True}, ensure_ascii=False))

    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
