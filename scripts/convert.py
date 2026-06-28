#!/usr/bin/env python
"""pptx を Markdown に変換する定型スクリプト。

markitdown を使ってスライド内容（見出し・本文・表・画像プレースホルダ・
スライド番号）を保持した Markdown を生成する。slide 数は python-pptx で数える。

実行には conda の py313 環境が必要:
    /opt/conda/envs/py313/bin/python src/pptx_to_md/convert.py ...

使い方:
    # pptx/ 配下の全 .pptx を変換
    python src/pptx_to_md/convert.py --all
    # 個別ファイルを指定
    python src/pptx_to_md/convert.py pptx/foo.pptx pptx/bar.pptx
    # 出力先を変更（既定: docs_from_agent）
    python src/pptx_to_md/convert.py --all --outdir docs_from_agent

結果は 1 ファイル 1 行の JSON で stdout に出力する（subagent が解析しやすいように）。
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from markitdown import MarkItDown
from pptx import Presentation

import pptx_images

# markitdown が出力する画像参照 ![alt](target) を捉える（alt と target を分離）。
_IMG_REF = re.compile(r"(!\[[^\]]*\]\()([^)]*)(\))")


def count_slides(pptx_path: Path) -> int | None:
    try:
        return len(Presentation(str(pptx_path)).slides)
    except Exception:
        return None


def rewrite_image_refs(text: str, targets: list[str | None]) -> tuple[str, list[str]]:
    """md 中の画像参照を、抽出画像への相対パスへ順番に置換する。

    targets は markitdown の画像参照出力順に並んだ相対パス（None はそのまま据え置き）。
    参照数と targets 数の不一致など、想定外があれば warnings に記録して続行する。
    """
    warnings: list[str] = []
    matches = list(_IMG_REF.finditer(text))
    if len(matches) != len(targets):
        warnings.append(
            f"image ref count ({len(matches)}) != extracted slot count ({len(targets)})"
        )

    it = iter(targets)

    def repl(m: re.Match) -> str:
        try:
            target = next(it)
        except StopIteration:
            return m.group(0)
        if target is None:
            return m.group(0)
        return f"{m.group(1)}images/{target}{m.group(3)}"

    return _IMG_REF.sub(repl, text), warnings


def convert_text(pptx_path: Path) -> str:
    return MarkItDown().convert(str(pptx_path)).text_content


def convert_one(pptx_path: Path, outdir: Path, overwrite: bool = True) -> dict:
    """1 つの pptx を Markdown に変換し、結果を dict で返す。"""
    result: dict = {"input": str(pptx_path), "output": None, "slides": None, "ok": False}

    if not pptx_path.is_file():
        result["error"] = "input file not found"
        return result

    # 出力は docs_from_agent/<名前>/<名前>.md にまとめる。
    out_path = outdir / pptx_path.stem / (pptx_path.stem + ".md")
    if out_path.exists() and not overwrite:
        result["error"] = "output exists (use --overwrite to replace)"
        result["output"] = str(out_path)
        return result

    try:
        text = convert_text(pptx_path)
    except Exception as exc:  # 変換失敗
        result["error"] = f"convert failed: {exc}"
        return result

    # 貼付画像を <名前>/images/ に抽出し、md の画像参照を実ファイルへ書き換える。
    images_dir = out_path.parent / "images"
    warnings: list[str] = []
    written: list[str] = []
    try:
        prs = Presentation(str(pptx_path))
        targets, written = pptx_images.extract_slide_images(prs, pptx_path.stem, images_dir)
        text, warnings = rewrite_image_refs(text, targets)
    except Exception as exc:
        warnings.append(f"image extraction failed: {exc}")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(text, encoding="utf-8")

    result.update(
        output=str(out_path),
        slides=count_slides(pptx_path),
        chars=len(text),
        images=len(written),
        images_dir=str(images_dir) if written else None,
        warnings=warnings,
        ok=True,
    )
    return result


def collect_inputs(args: argparse.Namespace) -> list[Path]:
    if args.all:
        return sorted(Path(args.pptx_dir).glob("**/*.pptx"))
    return [Path(p) for p in args.inputs]


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="pptx を Markdown に変換する")
    parser.add_argument("inputs", nargs="*", help="変換する .pptx ファイル")
    parser.add_argument("--all", action="store_true", help="--pptx-dir 配下の全 .pptx を変換")
    parser.add_argument("--pptx-dir", default="pptx", help="--all のときの探索元 (既定: pptx)")
    parser.add_argument("--outdir", default="docs_from_agent", help="出力先 (既定: docs_from_agent)")
    parser.add_argument("--overwrite", action="store_true", default=True, help="既存の .md を上書き (既定)")
    parser.add_argument("--no-overwrite", dest="overwrite", action="store_false", help="既存の .md を上書きしない")
    args = parser.parse_args(argv)

    inputs = collect_inputs(args)
    if not inputs:
        print("no .pptx files found", file=sys.stderr)
        return 1

    outdir = Path(args.outdir)
    failed = 0
    for pptx_path in inputs:
        res = convert_one(pptx_path, outdir, overwrite=args.overwrite)
        if not res["ok"]:
            failed += 1
        print(json.dumps(res, ensure_ascii=False))

    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
