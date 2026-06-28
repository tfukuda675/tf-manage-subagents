#!/usr/bin/env python
"""pptx に掲載されている表 (table) を CSV ファイルとして書き出すスクリプト。

各スライドの表（グループ内の表も含む）を探し、1 表 1 CSV として保存する。
日本語が Excel で文字化けしないよう UTF-8 (BOM 付き) で書き出す。

実行には conda の py313 環境が必要:
    /opt/conda/envs/py313/bin/python src/pptx_to_md/extract_tables.py --all

使い方:
    python src/pptx_to_md/extract_tables.py --all
    python src/pptx_to_md/extract_tables.py pptx/foo.pptx

結果は 1 ファイル 1 行の JSON で stdout に出力する。
"""
from __future__ import annotations

import argparse
import csv
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_tables(shapes):
    """図形コレクションから表 (graphicFrame の table) を再帰的に取り出す。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_tables(shape.shapes)
        elif getattr(shape, "has_table", False):
            yield shape.table


def table_to_rows(table) -> list[list[str]]:
    """表を文字列の二次元リストに変換する。セル内改行は空白に畳む。"""
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text.replace("\n", " ").strip() for cell in row.cells])
    return rows


def write_csv(rows: list[list[str]], path: Path) -> None:
    # Excel での文字化け回避のため UTF-8 BOM 付きで保存する。
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        csv.writer(f).writerows(rows)


def extract_one(pptx_path: Path, outroot: Path) -> dict:
    result: dict = {"input": str(pptx_path), "outdir": None, "tables": [], "count": 0, "ok": False}
    if not pptx_path.is_file():
        result["error"] = "input file not found"
        return result

    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        result["error"] = f"open failed: {exc}"
        return result

    out_dir = outroot / pptx_path.stem / "tables"
    # 再生成時に古い CSV が残らないよう掃除する。
    if out_dir.exists():
        for stale in out_dir.glob(f"{pptx_path.stem}_slide*_table*.csv"):
            stale.unlink()

    tables: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for tbl_idx, table in enumerate(iter_tables(slide.shapes), start=1):
            rows = table_to_rows(table)
            if not rows:
                continue
            out_dir.mkdir(parents=True, exist_ok=True)
            out = out_dir / f"{pptx_path.stem}_slide{slide_idx:03d}_table{tbl_idx}.csv"
            write_csv(rows, out)
            tables.append(str(out))

    result.update(
        outdir=str(out_dir) if tables else None,
        tables=tables,
        count=len(tables),
        ok=True,
    )
    return result


def collect_inputs(args: argparse.Namespace) -> list[Path]:
    if args.all:
        return sorted(Path(args.pptx_dir).glob("**/*.pptx"))
    return [Path(p) for p in args.inputs]


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="pptx の表を CSV に書き出す")
    parser.add_argument("inputs", nargs="*", help="対象の .pptx ファイル")
    parser.add_argument("--all", action="store_true", help="--pptx-dir 配下の全 .pptx を対象にする")
    parser.add_argument("--pptx-dir", default="pptx", help="--all のときの探索元 (既定: pptx)")
    parser.add_argument("--outdir", default="docs_from_agent", help="出力ルート (既定: docs_from_agent)")
    args = parser.parse_args(argv)

    inputs = collect_inputs(args)
    if not inputs:
        print("no .pptx files found", file=sys.stderr)
        return 1

    outroot = Path(args.outdir)
    outroot.mkdir(parents=True, exist_ok=True)
    failed = 0
    for pptx_path in inputs:
        res = extract_one(pptx_path, outroot)
        if not res["ok"]:
            failed += 1
        print(json.dumps(res, ensure_ascii=False))

    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
