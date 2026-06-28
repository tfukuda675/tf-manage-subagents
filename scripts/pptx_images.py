#!/usr/bin/env python
"""pptx 内の貼付画像 (Picture) を抽出する共有ロジック。

markitdown の pptx 変換と**同じ図形列挙・判定**を再現することで、convert.py が
生成する md の画像参照列と、抽出した画像ファイル列を 1:1 で対応づけられる。
convert.py（md の参照書き換え）と extract_images.py（単体抽出）の両方から使う。

参考: markitdown `_pptx_converter.py`
  - 画像参照は _is_picture が真の図形ごとに必ず出力（L149-152）
  - 反復順: スライド → shapes を (top,left) ソート → GROUP は再帰（L170-189）
  - 参照文字列 = re.sub(r"\\W","",shape.name)+".jpg"（L151）

実行には py313 環境 (python-pptx) と、EMF/WMF→PNG 変換のため soffice が必要。
"""
from __future__ import annotations

import re
import subprocess
import tempfile
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE

SOFFICE = "soffice"

# SVG 拡張 (svgBlip) を読むための名前空間。
_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
}
_R_EMBED = f"{{{_NS['r']}}}embed"

# md で表示できないため PNG に変換する形式。
_VECTOR_EXTS = {"emf", "wmf"}


def is_picture(shape) -> bool:
    """markitdown `_is_picture` と同一判定。"""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return hasattr(shape, "image")
    return False


def _sorted_shapes(shapes):
    """markitdown と同じ (top, left) 昇順ソート。"""
    return sorted(
        shapes,
        key=lambda x: (
            float("-inf") if not x.top else x.top,
            float("-inf") if not x.left else x.left,
        ),
    )


def iter_pictures_md_order(slide):
    """markitdown と同一順序で picture 図形を yield する。

    各スライドの shapes を (top,left) でソートし、GROUP は中身を同様にソートして
    再帰する。順序が markitdown の画像参照出力順と一致する。
    """

    def walk(shapes):
        for shape in _sorted_shapes(shapes):
            if is_picture(shape):
                yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes)

    yield from walk(slide.shapes)


def markitdown_ref(shape) -> str:
    """markitdown が md に出力する元の参照文字列を再現する（検証用）。"""
    return re.sub(r"\W", "", shape.name) + ".jpg"


def _svg_payload(shape):
    """picture が SVG 拡張を持てば (svg_blob, "svg") を返す。無ければ None。"""
    blip = shape._element.find(f".//{{{_NS['a']}}}blipFill/{{{_NS['a']}}}blip")
    if blip is None:
        return None
    svg = blip.find(f".//{{{_NS['asvg']}}}svgBlip")
    if svg is None:
        return None
    rid = svg.get(_R_EMBED)
    if not rid:
        return None
    try:
        part = shape.part.related_part(rid)
    except Exception:
        return None
    return part.blob, "svg"


def picture_payload(shape):
    """picture から (blob, ext) を取り出す。取得不能なら None。

    SVG 拡張があればベクタ (svg) を優先。それ以外は通常のラスタ blob。
    """
    svg = _svg_payload(shape)
    if svg is not None:
        return svg
    try:
        image = shape.image
        return image.blob, (image.ext or "img")
    except Exception:
        return None


def emf_wmf_to_png(src_path: Path, work_dir: Path) -> Path | None:
    """soffice で EMF/WMF を PNG に変換し、PNG パスを返す。失敗時 None。"""
    profile = work_dir / "lo_profile"
    try:
        subprocess.run(
            [
                SOFFICE,
                "--headless",
                f"-env:UserInstallation=file://{profile}",
                "--convert-to",
                "png",
                "--outdir",
                str(work_dir),
                str(src_path),
            ],
            check=True,
            capture_output=True,
        )
    except Exception:
        return None
    png = work_dir / (src_path.stem + ".png")
    return png if png.is_file() else None


def extract_slide_images(prs, stem: str, out_dir: Path) -> tuple[list[str | None], list[str]]:
    """全スライドの picture を markitdown 順で抽出する。

    戻り値 (targets, written):
      targets … markitdown の画像参照出力順にフラット化した、各 picture スロットの
                出力ファイルの「out_dir からの相対パス」。blob 取得不能スロットは None。
      written … 実際に書き出したファイルの絶対/相対パス一覧。
    命名は `<stem>_slideNNN_imgM.<ext>`（M はスライド内 picture スロット 1..）。
    ext が emf/wmf のときは PNG に変換し ext=png にする。
    """
    # 再生成時に古い画像を掃除。
    if out_dir.exists():
        for stale in out_dir.glob(f"{stem}_slide*_img*"):
            stale.unlink()

    targets: list[str | None] = []
    written: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slot = 0
        for shape in iter_pictures_md_order(slide):
            slot += 1
            payload = picture_payload(shape)
            if payload is None:
                targets.append(None)
                continue
            blob, ext = payload
            out_dir.mkdir(parents=True, exist_ok=True)
            base = f"{stem}_slide{slide_idx:03d}_img{slot}"

            if ext.lower() in _VECTOR_EXTS:
                # EMF/WMF は一旦書き出して soffice で PNG 変換する。
                with tempfile.TemporaryDirectory(dir=str(out_dir)) as tmp:
                    tmpdir = Path(tmp)
                    src = tmpdir / f"{base}.{ext.lower()}"
                    src.write_bytes(blob)
                    png = emf_wmf_to_png(src, tmpdir)
                    if png is None:
                        # 変換失敗時は元形式のまま保存（参照はするが表示はビューア依存）。
                        out = out_dir / f"{base}.{ext.lower()}"
                        out.write_bytes(blob)
                    else:
                        out = out_dir / f"{base}.png"
                        out.write_bytes(png.read_bytes())
            else:
                out = out_dir / f"{base}.{ext}"
                out.write_bytes(blob)

            targets.append(out.name)
            written.append(str(out))

    return targets, written
